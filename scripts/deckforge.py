#!/usr/bin/env python3
"""DeckForge local harness utilities.

This script provides deterministic glue for the DeckForge skills. It does not
replace Playwright MCP or frontend-design; it gives Codex repeatable local
commands for the file and render stages around those tools.
"""

from __future__ import annotations

import argparse
import json
import shutil
import subprocess
import sys
from pathlib import Path

from PIL import Image, ImageDraw, ImageOps, ImageStat
from pptx import Presentation
from pptx.util import Inches


SLIDE_W = 13.333333
SLIDE_H = 7.5


def fail(message: str) -> None:
    raise SystemExit(f"error: {message}")


def which(name: str) -> str:
    exe = shutil.which(name)
    if not exe:
        fail(f"missing required executable: {name}")
    return exe


def cmd_init(args: argparse.Namespace) -> None:
    root = Path(args.project).resolve()
    for child in ["captures", "slides", "renders", "pptx", "qa", "harness"]:
        (root / child).mkdir(parents=True, exist_ok=True)
    sources = root / "captures" / "sources.json"
    if not sources.exists():
        sources.write_text(json.dumps({"sources": []}, indent=2) + "\n", encoding="utf-8")
    print(root)


def cmd_html_to_png(args: argparse.Namespace) -> None:
    html = Path(args.html).resolve()
    if not html.exists():
        fail(f"missing HTML file: {html}")
    out_dir = Path(args.out_dir).resolve()
    out_dir.mkdir(parents=True, exist_ok=True)

    node = which("node")
    js = f"""
const {{ chromium }} = await import('playwright');
const browser = await chromium.launch({{ headless: true }});
const page = await browser.newPage({{ viewport: {{ width: {args.width}, height: {args.height} }}, deviceScaleFactor: 1 }});
await page.goto('file://{html}', {{ waitUntil: 'networkidle' }});
const slides = await page.locator('[data-slide]').count();
if (slides === 0) {{
  await page.screenshot({{ path: '{out_dir / "slide-01.png"}', fullPage: false }});
}} else {{
  for (let i = 0; i < slides; i++) {{
    const el = page.locator('[data-slide]').nth(i);
    await el.screenshot({{ path: `{out_dir}/slide-${{String(i + 1).padStart(2, '0')}}.png` }});
  }}
}}
await browser.close();
"""
    subprocess.run([node, "--input-type=module", "-e", js], check=True)
    print(f"rendered_dir: {out_dir}")


def image_files(path: Path, pattern: str) -> list[Path]:
    return [
        p for p in sorted(path.glob(pattern))
        if p.suffix.lower() in {".png", ".jpg", ".jpeg"}
    ]


def cmd_images_to_pptx(args: argparse.Namespace) -> None:
    input_dir = Path(args.input_dir).resolve()
    images = image_files(input_dir, args.pattern)
    if not images:
        fail(f"no images found in {input_dir} matching {args.pattern}")
    output = Path(args.output).resolve()
    output.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    blank = prs.slide_layouts[6]
    sw = float(prs.slide_width)
    sh = float(prs.slide_height)

    for idx, image in enumerate(images, start=1):
        slide = prs.slides.add_slide(blank)
        with Image.open(image) as im:
            iw, ih = im.size
        ratio = iw / ih
        slide_ratio = sw / sh
        if ratio > slide_ratio:
            h = sh
            w = h * ratio
            x = (sw - w) / 2
            y = 0
        else:
            w = sw
            h = w / ratio
            x = 0
            y = (sh - h) / 2
        pic = slide.shapes.add_picture(str(image), x, y, width=int(w), height=int(h))
        pic.name = f"DeckForge full-slide image {idx:02d}"

    prs.save(output)
    print(f"pptx: {output}")
    print(f"slides: {len(images)}")


def cmd_render_pptx(args: argparse.Namespace) -> None:
    pptx = Path(args.pptx).resolve()
    if not pptx.exists():
        fail(f"missing PPTX file: {pptx}")
    out_dir = Path(args.out_dir).resolve()
    out_dir.mkdir(parents=True, exist_ok=True)

    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice:
        fail("missing LibreOffice executable: soffice/libreoffice")
    pdftoppm = which("pdftoppm")

    subprocess.run([soffice, "--headless", "--convert-to", "pdf", "--outdir", str(out_dir), str(pptx)], check=True)
    pdf = out_dir / f"{pptx.stem}.pdf"
    if not pdf.exists():
        fail(f"expected PDF was not created: {pdf}")
    subprocess.run([pdftoppm, "-png", "-r", str(args.dpi), str(pdf), str(out_dir / "slide")], check=True)
    rendered = sorted(out_dir.glob("slide-*.png"))
    print(f"pdf: {pdf}")
    print(f"rendered: {len(rendered)}")


def cmd_qa(args: argparse.Namespace) -> None:
    render_dir = Path(args.render_dir).resolve()
    images = image_files(render_dir, args.pattern)
    if not images:
        fail(f"no render images found in {render_dir}")
    failures: list[str] = []
    for image in images:
        im = Image.open(image).convert("RGB")
        stat = ImageStat.Stat(im)
        variance = sum(stat.var) / 3
        bbox = ImageOps.invert(im.convert("L")).getbbox()
        print(f"{image.name}: size={im.size} variance={variance:.2f} nonwhite_bbox={bbox}")
        if bbox is None or variance < args.min_variance:
            failures.append(image.name)
    if args.contact_sheet:
        make_contact_sheet(images, Path(args.contact_sheet).resolve())
    if failures:
        print("failed:", ", ".join(failures))
        sys.exit(1)
    print(f"qa_ok: {len(images)}")


def make_contact_sheet(images: list[Path], output: Path) -> None:
    thumbs = []
    for image in images:
        im = Image.open(image).convert("RGB")
        im.thumbnail((320, 180))
        canvas = Image.new("RGB", (320, 205), "white")
        canvas.paste(im, ((320 - im.width) // 2, 0))
        ImageDraw.Draw(canvas).text((8, 184), image.name, fill=(0, 0, 0))
        thumbs.append(canvas)
    cols = 2
    rows = (len(thumbs) + cols - 1) // cols
    sheet = Image.new("RGB", (cols * 320, rows * 205), "#eeeeee")
    for idx, thumb in enumerate(thumbs):
        sheet.paste(thumb, ((idx % cols) * 320, (idx // cols) * 205))
    output.parent.mkdir(parents=True, exist_ok=True)
    sheet.save(output)
    print(f"contact_sheet: {output}")


def parser() -> argparse.ArgumentParser:
    p = argparse.ArgumentParser(description="DeckForge local harness")
    sub = p.add_subparsers(dest="cmd", required=True)

    init = sub.add_parser("init")
    init.add_argument("--project", default="deckforge")
    init.set_defaults(func=cmd_init)

    html = sub.add_parser("html-to-png")
    html.add_argument("--html", required=True)
    html.add_argument("--out-dir", required=True)
    html.add_argument("--width", type=int, default=1280)
    html.add_argument("--height", type=int, default=720)
    html.set_defaults(func=cmd_html_to_png)

    imgs = sub.add_parser("images-to-pptx")
    imgs.add_argument("--input-dir", required=True)
    imgs.add_argument("--pattern", default="slide-*.png")
    imgs.add_argument("--output", required=True)
    imgs.set_defaults(func=cmd_images_to_pptx)

    render = sub.add_parser("render-pptx")
    render.add_argument("--pptx", required=True)
    render.add_argument("--out-dir", required=True)
    render.add_argument("--dpi", type=int, default=120)
    render.set_defaults(func=cmd_render_pptx)

    qa = sub.add_parser("qa")
    qa.add_argument("--render-dir", required=True)
    qa.add_argument("--pattern", default="slide-*.png")
    qa.add_argument("--min-variance", type=float, default=1.0)
    qa.add_argument("--contact-sheet")
    qa.set_defaults(func=cmd_qa)

    return p


def main() -> None:
    args = parser().parse_args()
    args.func(args)


if __name__ == "__main__":
    main()
